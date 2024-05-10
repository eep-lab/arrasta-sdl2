import os
import pandas as pd
from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE

df = pd.DataFrame({
    'Fruit': ['Apples', 'Bananas', 'Cherries', 'Dates'],
    'Quantity': [50, 30, 20, 10]
})

prs = Presentation()
slide = prs.slides.add_slide(prs.slide_layouts[0])
for shape in slide.placeholders:
    print('%d %s' % (shape.placeholder_format.idx, shape.name))

title = slide.shapes.title
subtitle = slide.placeholders[1]
try:
    subtitle = slide.placeholders[1]
    print(subtitle.name)
except KeyError:
    print("Placeholder not found with idx %d" % subtitle.idx)

title.text = "Hello, Python-powered PowerPoint!"
subtitle.text = "Creating presentations with python-pptx."

slide = prs.slides.add_slide(prs.slide_layouts[5])


chart_data = CategoryChartData()
chart_data.categories = df['Fruit']
chart_data.add_series('Series 1', df['Quantity'])

# x, y, cx, cy = 2, 2, 6, 4.5
# chart = slide.shapes.add_chart(
#     XL_CHART_TYPE.BAR_CLUSTERED, x, y, cx, cy, chart_data).chart


# add chart to slide
x, y, cx, cy = 1, 1, 5, 5
chart = slide.shapes.add_chart(
    XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data
).chart


output = os.path.join(os.getcwd(), 'output', 'pandas-python-pptx.pptx')
prs.save(output)